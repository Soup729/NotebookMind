import argparse
import json
from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.util import Pt
from reportlab.lib.pagesizes import A4
from reportlab.pdfbase.pdfmetrics import stringWidth
from reportlab.pdfgen import canvas


def load_payload(path: str) -> dict:
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


def ensure_parent(path: str) -> None:
    Path(path).parent.mkdir(parents=True, exist_ok=True)


def render_docx(payload: dict, output_path: str) -> None:
    document = Document()
    document.add_heading(payload.get("title", "Notebook Export"), 0)
    requirements = payload.get("requirements", "").strip()
    if requirements:
        document.add_paragraph(requirements)
    for section in payload.get("outline", []):
        document.add_heading(section.get("heading", "Section"), level=1)
        for bullet in section.get("bullets", []):
            document.add_paragraph(str(bullet), style="List Bullet")
    refs = payload.get("source_refs", [])
    if refs:
        document.add_heading("Sources", level=1)
        for ref in refs:
            line = ref.get("document_name", "Unknown document")
            page = ref.get("page")
            quote = ref.get("quote", "").strip()
            if page:
                line += f", Page {page}"
            if quote:
                line += f": {quote}"
            document.add_paragraph(line, style="List Bullet")
    document.save(output_path)


def render_pptx(payload: dict, output_path: str) -> None:
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = payload.get("title", "Notebook Export")
    title_slide.placeholders[1].text = payload.get("requirements", "") or "PPTX export"

    for section in payload.get("outline", []):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = section.get("heading", "Section")
        text_frame = slide.shapes.placeholders[1].text_frame
        text_frame.clear()
        for idx, bullet in enumerate(section.get("bullets", [])):
            paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
            paragraph.text = str(bullet)
            paragraph.level = 0
            paragraph.font.size = Pt(20)

    refs = payload.get("source_refs", [])
    if refs:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Sources"
        text_frame = slide.shapes.placeholders[1].text_frame
        text_frame.clear()
        for idx, ref in enumerate(refs):
            line = ref.get("document_name", "Unknown document")
            page = ref.get("page")
            quote = ref.get("quote", "").strip()
            if page:
                line += f", Page {page}"
            if quote:
                line += f": {quote}"
            paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
            paragraph.text = line
            paragraph.level = 0
            paragraph.font.size = Pt(14)
    prs.save(output_path)


def wrap_lines(text: str, font_name: str, font_size: int, max_width: float):
    words = text.split()
    if not words:
        return [""]
    lines = []
    current = words[0]
    for word in words[1:]:
        candidate = current + " " + word
        if stringWidth(candidate, font_name, font_size) <= max_width:
            current = candidate
        else:
            lines.append(current)
            current = word
    lines.append(current)
    return lines


def render_pdf(payload: dict, output_path: str) -> None:
    c = canvas.Canvas(output_path, pagesize=A4)
    width, height = A4
    margin = 50
    y = height - margin

    def new_page():
        nonlocal y
        c.showPage()
        y = height - margin

    def draw_text_block(text: str, font_name: str, font_size: int, indent: int = 0, gap: int = 4):
        nonlocal y
        max_width = width - margin * 2 - indent
        for line in wrap_lines(text, font_name, font_size, max_width):
            if y < margin:
                new_page()
            c.setFont(font_name, font_size)
            c.drawString(margin + indent, y, line)
            y -= font_size + gap

    draw_text_block(payload.get("title", "Notebook Export"), "Helvetica-Bold", 18, 0, 8)
    requirements = payload.get("requirements", "").strip()
    if requirements:
        draw_text_block(requirements, "Helvetica", 11, 0, 6)
    for section in payload.get("outline", []):
        draw_text_block(section.get("heading", "Section"), "Helvetica-Bold", 14, 0, 6)
        for bullet in section.get("bullets", []):
            draw_text_block("- " + str(bullet), "Helvetica", 11, 16, 4)
        y -= 6
    refs = payload.get("source_refs", [])
    if refs:
        draw_text_block("Sources", "Helvetica-Bold", 14, 0, 6)
        for ref in refs:
            line = ref.get("document_name", "Unknown document")
            page = ref.get("page")
            quote = ref.get("quote", "").strip()
            if page:
                line += f", Page {page}"
            if quote:
                line += f": {quote}"
            draw_text_block("- " + line, "Helvetica", 10, 16, 4)
    c.save()


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    payload = load_payload(args.input)
    ensure_parent(args.output)
    fmt = str(payload.get("format", "")).strip().lower()

    if fmt == "docx":
        render_docx(payload, args.output)
    elif fmt == "pptx":
        render_pptx(payload, args.output)
    elif fmt == "pdf":
        render_pdf(payload, args.output)
    else:
        raise SystemExit(f"unsupported format: {fmt}")

    print(json.dumps({"ok": True, "output": args.output}, ensure_ascii=False))


if __name__ == "__main__":
    main()
